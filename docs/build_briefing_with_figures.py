#!/usr/bin/env python3
"""
Augment the PI briefing deck with the census data-analysis figures.
===================================================================

Opens the 2026-06-26 briefing, inserts the figures from
paper/replication/verification/figures/ immediately after their related
text slides, and writes docs/VERIFIER_CENSUS_BRIEFING_2026-06-27.pptx.

Run:
  .venv-django/bin/python docs/build_briefing_with_figures.py
"""
from __future__ import annotations

import re
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "docs/VERIFIER_CENSUS_BRIEFING_2026-06-26.pptx"
DST = ROOT / "docs/VERIFIER_CENSUS_BRIEFING_2026-06-27.pptx"
FIG = ROOT / "paper/replication/verification/figures"

# figure -> (title, caption, normalized anchor substring of the slide it follows)
PLAN = [
    ("fig1_corpus_funnel.png",
     "Corpus funnel (10,200 → 341)",
     "10,200 enumerated → 10,103 with body → 1,939 with a test claim → 341 with a checkable claim.",
     "built (census)"),
    ("fig2_headline_outcome.png",
     "Outcome over 3,005 checkable claims",
     "Consistent vs inconsistent (333, 11.1%) vs decision-changing (52, 1.7%).",
     "Census headline"),
    ("fig4_reported_vs_recomputed_p.png",
     "Reported vs recomputed p (333 flags)",
     "Log-log; ★ = decision-changing; colored by false-positive category.",
     "Census headline"),
    ("fig3_fp_validation.png",
     "False-positive validation of the 333 flags",
     "TRUE_LIKELY 262 · REVIEW 25 · one-tailed 46 · mis-extraction 0 (was 157).",
     "PRE-FIX"),
    ("fig6_rate_robustness.png",
     "The rate is single-digit across every frame",
     "Raw 11.1% · IPW 10.5% · likely-true 8.7% · independent OA 5.6%.",
     "Robustness"),
    ("fig5_by_statistic_type.png",
     "Flagged inconsistencies by statistic type",
     "t / F / r / z / chi-square — all flagged vs likely-true.",
     "Robustness"),
]

DARK = RGBColor(0x21, 0x21, 0x21)
GREY = RGBColor(0x55, 0x55, 0x55)


def norm(s: str) -> str:
    return re.sub(r"\s+", " ", s.replace("\xa0", " ")).strip()


def slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            parts.append(sh.text_frame.text)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return norm(" ".join(parts))


def make_fig_slide(prs, title, caption, img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    sw, sh = prs.slide_width, prs.slide_height
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.22), sw - Inches(0.9), Inches(0.7))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = DARK
    cb = slide.shapes.add_textbox(Inches(0.45), sh - Inches(0.7), sw - Inches(0.9), Inches(0.5))
    cp = cb.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = caption
    cr.font.size = Pt(12); cr.font.color.rgb = GREY; cr.font.italic = True
    pic = slide.shapes.add_picture(str(img), 0, 0)
    avail_w, avail_h = sw - Inches(1.0), sh - Inches(1.9)
    scale = min(avail_w / pic.width, avail_h / pic.height)
    pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
    pic.left = int((sw - pic.width) / 2)
    pic.top = Inches(1.05) + int((avail_h - pic.height) / 2)
    return slide


def main() -> int:
    prs = Presentation(str(SRC))
    n_orig = len(prs.slides._sldIdLst)
    orig_texts = [slide_text(s) for s in prs.slides]

    # add figure slides (they append to the end), aligned with PLAN order
    fig_specs = []
    for fname, title, caption, anchor in PLAN:
        img = FIG / fname
        if not img.exists():
            print(f"  !! missing {img}; skipping")
            continue
        make_fig_slide(prs, title, caption, img)
        fig_specs.append((fname, anchor))

    lst = prs.slides._sldIdLst
    all_nodes = list(lst)
    orig_nodes, fig_nodes = all_nodes[:n_orig], all_nodes[n_orig:]

    # map each figure to the ORIGINAL slide index it should follow
    after = defaultdict(list)
    for (fname, anchor), fnode in zip(fig_specs, fig_nodes):
        na = norm(anchor)
        idx = max((i for i, t in enumerate(orig_texts) if na in t), default=n_orig - 1)
        after[idx].append((fname, fnode))
        print(f"  + {fname}  -> after slide {idx + 1}")

    # rebuild order: each original slide, then its figures (PLAN order)
    desired = []
    for i, onode in enumerate(orig_nodes):
        desired.append(onode)
        desired.extend(fn for _, fn in after.get(i, []))

    for node in list(lst):
        lst.remove(node)
    for node in desired:
        lst.append(node)

    prs.save(str(DST))
    print(f"\nSaved {DST}  ({len(prs.slides._sldIdLst)} slides, was {n_orig})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
